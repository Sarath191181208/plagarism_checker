from pptx import Presentation 
from .base_class import BaseParser

class PptParser(BaseParser):
    @staticmethod
    def parse(abs_file_name: str) -> str:
            
               prs = Presentation(abs_file_name)
               text_runs=[]
               for slide in prs.slides:
                   for shape in slide.shape:
                      if not shape.has_text_frame:
                         continue
                      for paragraph in shape.text_frame.paragraphs:
                          for run in paragraph.runs:
                               text_runs.append(run.text)
               return '\n'.join(text_runs)

                
                
                



    
